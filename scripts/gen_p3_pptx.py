from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

BLACK = RGBColor(0, 0, 0)

# 生成模块框
def add_box(text, x, y, w, h, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = BLACK
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = BLACK
    return shape

# 生成信号文字标签
def add_label(text, x, y, size=11, w=2.2):
    box = slide.shapes.add_textbox(x, y, Inches(w), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = BLACK
    return box

# 生成连接器并绑定两端
def add_conn(connector_type, begin_shape, end_shape, begin_idx=1, end_idx=3):
    conn = slide.shapes.add_connector(connector_type, 0, 0, 0, 0)
    conn.line.color.rgb = BLACK
    conn.begin_connect(begin_shape, begin_idx)
    conn.end_connect(end_shape, end_idx)
    return conn

# 主模块（左->右）
regfile = add_box("Register File\n(regfile8x16)", Inches(0.6), Inches(2.1), Inches(2.2), Inches(1.2))
alu = add_box("ALU\n(alu16)", Inches(3.6), Inches(2.2), Inches(1.6), Inches(1.0))
mem = add_box("Data Memory\n(data_mem32x16)", Inches(6.0), Inches(2.1), Inches(2.2), Inches(1.2))

# MUX 使用梯形表示
wb_mux = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(8.6), Inches(2.2), Inches(1.4), Inches(1.0))
wb_mux.fill.solid()
wb_mux.fill.fore_color.rgb = RGBColor(255, 255, 255)
wb_mux.line.color.rgb = BLACK
wb_tf = wb_mux.text_frame
wb_tf.clear()
wb_p = wb_tf.paragraphs[0]
wb_run = wb_p.add_run()
wb_run.text = "WB MUX\n(mem_to_reg)"
wb_run.font.size = Pt(11)
wb_run.font.color.rgb = BLACK

# 辅助模块（ALU 第二操作数选择）
alu_src_mux = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(2.6), Inches(0.9), Inches(1.6), Inches(0.8))
alu_src_mux.fill.solid()
alu_src_mux.fill.fore_color.rgb = RGBColor(255, 255, 255)
alu_src_mux.line.color.rgb = BLACK
alu_tf = alu_src_mux.text_frame
alu_tf.clear()
alu_p = alu_tf.paragraphs[0]
alu_run = alu_p.add_run()
alu_run.text = "ALU Src MUX\n(alu_src_imm)"
alu_run.font.size = Pt(10)
alu_run.font.color.rgb = BLACK

# 数据通路连接（水平/垂直+90度，端点绑定模块）
add_conn(MSO_CONNECTOR.ELBOW, regfile, alu, begin_idx=1, end_idx=3)
add_conn(MSO_CONNECTOR.ELBOW, regfile, alu_src_mux, begin_idx=1, end_idx=3)
add_conn(MSO_CONNECTOR.ELBOW, alu_src_mux, alu, begin_idx=1, end_idx=0)
add_conn(MSO_CONNECTOR.ELBOW, alu, mem, begin_idx=1, end_idx=3)
add_conn(MSO_CONNECTOR.ELBOW, mem, wb_mux, begin_idx=1, end_idx=3)

# ALU 结果回写选择路径（绕过 mem）
add_conn(MSO_CONNECTOR.ELBOW, alu, wb_mux, begin_idx=1, end_idx=1)

# 写回到寄存器堆
add_conn(MSO_CONNECTOR.ELBOW, wb_mux, regfile, begin_idx=3, end_idx=1)

# 端口与控制信号标签（顶部留白区）
add_label("rs1[2:0]", Inches(0.2), Inches(2.0))
add_label("rs2[2:0]", Inches(0.2), Inches(2.4))
add_label("rd[2:0]", Inches(0.2), Inches(2.8))
add_label("reg_we", Inches(0.2), Inches(3.2))
add_label("init_we/init_addr/init_data", Inches(0.2), Inches(3.6), size=10, w=3.0)

add_label("imm[15:0]", Inches(2.0), Inches(0.6))
add_label("alu_op[2:0]", Inches(3.4), Inches(1.6))
add_label("mem_we", Inches(5.6), Inches(1.6))
add_label("byte_en[1:0]", Inches(5.6), Inches(1.9))
add_label("mem_to_reg", Inches(8.0), Inches(1.6))
add_label("alu_src_imm", Inches(2.6), Inches(0.3), size=10)

# 数据线标注
add_label("src_a", Inches(2.9), Inches(2.05), size=10)
add_label("src_b", Inches(2.9), Inches(2.35), size=10)
add_label("alu_result[15:0]", Inches(4.6), Inches(2.0), size=10, w=2.0)
add_label("addr = alu_result[4:0]", Inches(6.0), Inches(3.4), size=9, w=2.8)
add_label("mem_rdata[15:0]", Inches(7.5), Inches(2.0), size=10, w=2.0)
add_label("wb_data[15:0]", Inches(8.2), Inches(3.4), size=10, w=1.8)

out_path = "d:/Program/数电实践/assignments/1/p3/docs/p3_datapath.pptx"
prs.save(out_path)
print(out_path)
